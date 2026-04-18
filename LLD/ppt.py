from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation with custom poster size (13x7.5 inches)
prs = Presentation()
prs.slide_width = Inches(13)
prs.slide_height = Inches(7.5)

# Add blank slide layout
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# Set background color (dark navy blue)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 32, 64)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "EFFECTIVE AI TOOLS EVERY STUDENT MUST KNOW"
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Subtitle / name
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.5))
sub_frame = sub_box.text_frame
sub_frame.text = "Your Name | National Textile University, Faisalabad"
sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
sub_frame.paragraphs[0].font.size = Pt(18)
sub_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

# Column 1: Text & Study Tools
col1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(3.8), Inches(4.5))
tf1 = col1.text_frame
tf1.text = "📝 TEXT & STUDY TOOLS\n\n"
tf1.paragraphs[0].font.size = Pt(22)
tf1.paragraphs[0].font.bold = True
tf1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)
for tool, desc in [("ChatGPT", "Essays, summaries, ideas"),
                   ("Google Gemini", "Research + images"),
                   ("Perplexity AI", "Cited answers"),
                   ("Elicit", "Academic papers"),
                   ("ChatPDF", "Talk to PDFs")]:
    p = tf1.add_paragraph()
    p.text = f"🔹 {tool}\n   {desc}\n"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)

# Column 2: Design & Productivity
col2 = slide.shapes.add_textbox(Inches(4.6), Inches(1.8), Inches(3.8), Inches(4.5))
tf2 = col2.text_frame
tf2.text = "🎨 DESIGN & PRODUCTIVITY\n\n"
tf2.paragraphs[0].font.size = Pt(22)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)
for tool, desc in [("Canva AI", "Posters, presentations"),
                   ("Gamma.ai", "Generate PPT slides"),
                   ("Beautiful.ai", "Auto-design"),
                   ("Notion AI", "Smart notes"),
                   ("Otter.ai", "Lecture transcription")]:
    p = tf2.add_paragraph()
    p.text = f"🔹 {tool}\n   {desc}\n"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)

# Column 3: Coding, Math & Responsible Use
col3 = slide.shapes.add_textbox(Inches(8.7), Inches(1.8), Inches(3.8), Inches(4.5))
tf3 = col3.text_frame
tf3.text = "💻 CODING & MATH\n\n"
tf3.paragraphs[0].font.size = Pt(22)
tf3.paragraphs[0].font.bold = True
tf3.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)
for tool, desc in [("GitHub Copilot", "Code completion"),
                   ("Codeium", "Free AI coding"),
                   ("Wolfram Alpha", "Math & engineering"),
                   ("Photomath", "Solve by camera")]:
    p = tf3.add_paragraph()
    p.text = f"🔹 {tool}\n   {desc}\n"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)

# Add warning / responsible use box
warning = slide.shapes.add_textbox(Inches(8.7), Inches(4.2), Inches(3.8), Inches(2.0))
wf = warning.text_frame
wf.text = "⚠️ RESPONSIBLE USE\n\n"
wf.paragraphs[0].font.size = Pt(18)
wf.paragraphs[0].font.bold = True
wf.paragraphs[0].font.color.rgb = RGBColor(255, 100, 100)
for line in ["Verify facts", "No copy-paste", "Respect privacy"]:
    p = wf.add_paragraph()
    p.text = f"• {line}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)

# Bottom banner
banner = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
banner.fill.solid()
banner.fill.fore_color.rgb = RGBColor(0, 80, 120)
banner.line.color.rgb = RGBColor(255, 255, 0)
bf = banner.text_frame
bf.text = "🔐 Stay Smart, Work Faster – AI is Your Assistant, Not a Shortcut"
bf.paragraphs[0].alignment = PP_ALIGN.CENTER
bf.paragraphs[0].font.size = Pt(20)
bf.paragraphs[0].font.bold = True
bf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Save file
prs.save("AI_Tools_Poster.pptx")
print("✅ Poster created: AI_Tools_Poster.pptx")